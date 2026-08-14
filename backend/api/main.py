"""FastAPI app — the real ingestion + query pipeline behind the React UI.

Endpoints (match frontend/src/api.jsx):
    POST   /upload         multipart file -> run ingestion, return doc metadata
    POST   /files/stage    multipart file -> save only (no ingest) -> {file_path}, for chat attach
    GET    /files          list ingested documents
    GET    /files/{id}     one document's metadata
    DELETE /files/{id}     delete a document (chunks cascade)
    POST   /agent/chat     {message, session_id?, approved_writes?} -> agent picks a tool
                           (ingest_document / search_documents / list_documents / sql_read);
                           writes need approval
    GET    /agent/sessions           list past agent-chat conversations (sidebar)
    GET    /agent/sessions/{id}      one conversation's full turn history
    DELETE /agent/sessions/{id}      delete a conversation
    GET    /files/{id}/original      raw bytes of an image document, for direct <img> viewing
    GET    /files/{id}/docx-html     docx converted to HTML, for in-panel viewing
    GET    /health         liveness

Everything document-facing goes through the agent now — there is no direct
(non-agentic) RAG endpoint. search_documents is just another tool the agent
calls; it is not exposed as its own HTTP route.

This wires the WHOLE pipeline (categorize -> extract -> ... -> index for ingest;
plan -> retrieve -> answer for chat), not just categorization. Running it needs
the stack up (Postgres, Qdrant) and models available — see docker-compose + .env.

Run:  uvicorn backend.api.main:app --reload --port 8000
"""
from __future__ import annotations

import glob
import logging
import mimetypes
import os
import shutil
import uuid

# ── Logging setup (before any imports that may emit logs) ─────────────────────
logging.basicConfig(
    level=logging.INFO,
    format="%(message)s",
)
# Silence chatty third-party and internal libraries that add noise on startup
_noisy_loggers = [
    "httpx", "httpcore", "openai", "anthropic", "hpack", "h2",
    "urllib3", "requests", "PIL", "fastembed",
    "backend.pipeline.default_registry",
    "backend.guardrails.startup_check",
    "backend.core.health_probe"
]
for _logger_name in _noisy_loggers:
    logging.getLogger(_logger_name).setLevel(logging.WARNING)

class EndpointFilter(logging.Filter):
    """Filter out noisy background polling endpoints (/progress, /health, /pages/.../image, /pdf-info) from Uvicorn access logs."""
    def filter(self, record: logging.LogRecord) -> bool:
        msg = record.getMessage()
        return not (
            "/progress" in msg
            or "/health" in msg
            or "/pages/" in msg
            or "/image" in msg
            or "/pdf-info" in msg
        )

logging.getLogger("uvicorn.access").addFilter(EndpointFilter())

from backend.core.db_logging import setup_db_logging
# We don't hold the reference to listener here, it runs natively in background
setup_db_logging(level=logging.INFO)




import yaml
from dotenv import load_dotenv

# Load .env BEFORE load_config so ${GROQ_API_KEY}/${POSTGRES_URL}/... resolve.
load_dotenv()

from fastapi import BackgroundTasks, FastAPI, File, HTTPException, UploadFile, Response, Request  # noqa: E402
from fastapi.middleware.cors import CORSMiddleware  # noqa: E402
from fastapi.staticfiles import StaticFiles  # noqa: E402
from pydantic import BaseModel  # noqa: E402

from backend.agent.executor import run_agent  # noqa: E402
from backend.agent_tools import build_agent_registry  # noqa: E402
from backend.core.config import load_config  # noqa: E402
from backend.core.yaml_handler import (  # noqa: E402
    load_yaml_roundtrip,
    dump_yaml_roundtrip,
    apply_settings_in_place,
)
from backend.core.models import warm_up  # noqa: E402
from backend.core.llm_client import clean_message_content  # noqa: E402
from backend.pipeline.default_registry import build_default_registry  # noqa: E402
from backend.pipeline.ingest import ingest_document  # noqa: E402
from backend.storage.postgres_store import PostgresStore  # noqa: E402
from backend.storage.qdrant_store import QdrantStore  # noqa: E402
from langchain_core.messages import AIMessage, HumanMessage  # noqa: E402

from backend.guardrails.token_quota import get_enforcer, get_reserve_tokens
from backend.guardrails.startup_check import run_startup_self_test
from backend.core.health_probe import background_health_loop
from backend.storage.postgres_store import dsn_from_env as _pg_dsn


logger = logging.getLogger(__name__)

CONFIG_PATH = os.getenv("CONFIG_PATH", "config/global.yaml")
UPLOAD_DIR = os.getenv("UPLOAD_DIR", "uploads")

EXT_TO_FILE_TYPE = {
    ".pdf": "pdf",
    ".xlsx": "excel", ".xls": "excel", ".xlsm": "excel",
    ".pptx": "ppt", ".ppt": "ppt",
    ".docx": "docx", ".doc": "docx",
    ".png": "image", ".jpg": "image", ".jpeg": "image", ".tif": "image", ".tiff": "image",
}

app = FastAPI(title="Document Intelligence + RAG Accelerator", version="1.0.0")


@app.on_event("startup")
def startup_event():
    pass





# CORS origins: configurable via ALLOWED_ORIGINS env var (comma-separated).
# Default covers local dev; set ALLOWED_ORIGINS in production .env.
_ALLOWED_ORIGINS = [
    o.strip()
    for o in os.getenv(
        "ALLOWED_ORIGINS",
        "http://localhost:5173,http://localhost:3000,http://127.0.0.1:5173,http://127.0.0.1:3000",
    ).split(",")
    if o.strip()
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=_ALLOWED_ORIGINS,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Serve enriched images so citation image_path (/images/<doc>/<file>) resolves.
# vision writes them under uploads/images/<doc_id>/<block_id>.png.
_IMAGES_DIR = os.path.join(UPLOAD_DIR, "images")
os.makedirs(_IMAGES_DIR, exist_ok=True)
app.mount("/images", StaticFiles(directory=_IMAGES_DIR), name="images")

# Serve full-page images (visual grounding). Rendered at ingest for PDF pages that
# produced chunks; stored under uploads/pages/<doc_id>/p{N}.jpg, served at /pages/...
_PAGES_DIR = os.path.join(UPLOAD_DIR, "pages")
os.makedirs(_PAGES_DIR, exist_ok=True)
app.mount("/pages", StaticFiles(directory=_PAGES_DIR), name="pages")


def get_unique_path(directory: str, filename: str) -> str:
    base, ext = os.path.splitext(filename)
    counter = 1
    unique_path = os.path.join(directory, filename)
    while os.path.exists(unique_path):
        unique_path = os.path.join(directory, f"{base}_{counter}{ext}")
        counter += 1
    return unique_path


def auto_ingestion_loop() -> None:
    import time
    logger.info("Auto-ingestion watcher thread loop started.")
    while True:
        try:
            # Note: _config is reloaded in-place by _reload_pipeline()
            cfg = _config.get("auto_ingestion") or {}
            enabled = cfg.get("enabled", False)
            
            if enabled:
                watch_dir_raw = cfg.get("watch_dir", "auto_ingest")
                poll_interval = cfg.get("poll_interval", 10)
                on_success = cfg.get("on_success", "move")
                on_failure = cfg.get("on_failure", "move")
                
                # Resolve watch_dir relative to the project root
                base_dir = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
                watch_dir = os.path.abspath(watch_dir_raw) if os.path.isabs(watch_dir_raw) else os.path.abspath(os.path.join(base_dir, watch_dir_raw))
                
                if not os.path.exists(watch_dir):
                    try:
                        os.makedirs(watch_dir, exist_ok=True)
                        logger.info("Created watch directory: %s", watch_dir)
                    except Exception as e:
                        logger.error("Failed to create watch directory %s: %s", watch_dir, e)
                        time.sleep(poll_interval)
                        continue
                
                # Ensure destination folders for move operations exist
                processed_dir = os.path.join(watch_dir, "processed")
                failed_dir = os.path.join(watch_dir, "failed")
                if on_success == "move":
                    os.makedirs(processed_dir, exist_ok=True)
                if on_failure == "move":
                    os.makedirs(failed_dir, exist_ok=True)
                
                # Scan watch directory for files
                for item in sorted(os.listdir(watch_dir)):
                    src_path = os.path.join(watch_dir, item)
                    
                    # Skip subdirectories (like processed/failed)
                    if os.path.isdir(src_path):
                        continue
                    
                    # Skip hidden files
                    if item.startswith("."):
                        continue
                    
                    # Wait for copy to complete (file size must be stable)
                    try:
                        initial_size = os.path.getsize(src_path)
                        time.sleep(2)
                        current_size = os.path.getsize(src_path)
                        if initial_size != current_size or current_size == 0:
                            continue
                    except Exception:
                        continue  # File might have been renamed or deleted
                    
                    # Map to known file type
                    file_type = _file_type(item)
                    if file_type == "unknown":
                        logger.warning("Auto-ingestion: Unsupported file type for %s. Moving/Deleting file.", item)
                        if on_failure == "move":
                            try:
                                shutil.move(src_path, get_unique_path(failed_dir, item))
                            except Exception as e:
                                logger.error("Auto-ingestion: Failed to move unsupported file %s: %s", item, e)
                        else:
                            try:
                                os.remove(src_path)
                            except Exception as e:
                                logger.error("Auto-ingestion: Failed to remove unsupported file %s: %s", item, e)
                        continue
                    
                    logger.info("Auto-ingestion: Processing file %s...", item)
                    document_id = str(uuid.uuid4())
                    logger.info("Auto-ingestion: Starting ingestion of '%s' (doc: %s)", item, document_id)
                    t0 = time.time()
                    
                    dest = os.path.join(UPLOAD_DIR, f"{document_id}_{item}")
                    os.makedirs(UPLOAD_DIR, exist_ok=True)
                    
                    try:
                        shutil.copy2(src_path, dest)
                    except Exception as e:
                        logger.error("Auto-ingestion: Failed to copy %s to upload dir: %s", item, e)
                        print(f"=== Auto-Ingestion: Failed to copy '{item}' to upload dir ===\n", flush=True)
                        continue
                    
                    # Insert record into DB as processing
                    pg = PostgresStore()
                    try:
                        pg.insert_document(document_id, item, file_type, dest)
                    except Exception as e:
                        logger.error("Auto-ingestion: Failed to insert document record for %s: %s", item, e)
                        print(f"=== Auto-Ingestion: Failed to record database entry for '{item}' ===\n", flush=True)
                        pg.close()
                        if os.path.exists(dest):
                            try:
                                os.remove(dest)
                            except Exception:
                                pass
                        continue
                    finally:
                        pg.close()
                    
                    # Run the ingestion synchronously in this thread
                    try:
                        _run_ingestion(document_id, dest, file_type, item)
                    except Exception as e:
                        logger.error("Auto-ingestion: Exception while ingesting %s: %s", item, e)
                    
                    # Check final ingestion status from database
                    status = "failed"
                    pg = PostgresStore()
                    try:
                        doc = pg.get_document(document_id)
                        if doc:
                            status = doc.get("status", "failed")
                    except Exception as e:
                        logger.error("Auto-ingestion: Failed to check final status of %s: %s", item, e)
                    finally:
                        pg.close()
                    
                    elapsed = time.time() - t0
                    if status == "ready":
                        logger.info("Auto-ingestion: Ingested %s successfully.", item)
                        print(f"=== Auto-Ingestion: Ingested '{item}' successfully in {elapsed:.1f}s ===\n", flush=True)
                        if on_success == "move":
                            try:
                                shutil.move(src_path, get_unique_path(processed_dir, item))
                            except Exception as e:
                                logger.error("Auto-ingestion: Failed to move successful file %s: %s", item, e)
                        else:
                            try:
                                os.remove(src_path)
                            except Exception as e:
                                logger.error("Auto-ingestion: Failed to delete successful file %s: %s", item, e)
                    else:
                        logger.error("Auto-ingestion: Failed to ingest %s (status=%s).", item, status)
                        print(f"=== Auto-Ingestion: Failed to ingest '{item}' in {elapsed:.1f}s (status: {status}) ===\n", flush=True)
                        if on_failure == "move":
                            try:
                                shutil.move(src_path, get_unique_path(failed_dir, item))
                            except Exception as e:
                                logger.error("Auto-ingestion: Failed to move failed file %s: %s", item, e)
                        else:
                            try:
                                os.remove(src_path)
                            except Exception as e:
                                logger.error("Auto-ingestion: Failed to delete failed file %s: %s", item, e)
            
        except Exception as e:
            logger.exception("Error in auto-ingestion loop: %s", e)
        
        # Determine poll interval from config or default to 10
        try:
            cfg = _config.get("auto_ingestion") or {}
            poll_interval = int(cfg.get("poll_interval", 10))
        except Exception:
            poll_interval = 10
            
        time.sleep(max(1, poll_interval))


def start_auto_ingestion_watcher() -> None:
    import threading
    t = threading.Thread(target=auto_ingestion_loop, name="AutoIngestionWatcher", daemon=True)
    t.start()


@app.on_event("startup")
async def on_startup():
    """Verify connections and auto-initialize database schema if missing."""
    import psycopg
    import re
    import asyncio
    from qdrant_client import QdrantClient
    
    from backend.storage.postgres_store import dsn_from_env
    postgres_url = dsn_from_env()
    qdrant_url = os.getenv("QDRANT_URL", "http://localhost:6333")
    qdrant_api_key = os.getenv("QDRANT_API_KEY")
    
    logger.info("Verifying database and vector store connections...")
    
    # 1. Verify / Initialize Postgres
    if postgres_url:
        try:
            conn = psycopg.connect(postgres_url, connect_timeout=5, prepare_threshold=None)
            with conn.cursor() as cur:
                # Run init_db.sql idempotently to ensure all tables and indexes (including guardrails) exist
                logger.debug("Syncing database schema using scripts/init_db.sql...")
                base_dir = os.path.dirname(os.path.dirname(os.path.dirname(__file__)))
                schema_path = os.path.join(base_dir, "scripts", "init_db.sql")
                if os.path.exists(schema_path):
                    with open(schema_path, "r", encoding="utf-8") as sf:
                        sql_content = sf.read()
                    
                    sql_no_comments = re.sub(r"--[^\n]*", "", sql_content)
                    statements = []
                    for stmt in sql_no_comments.split(";"):
                        stmt = stmt.strip()
                        if stmt and not stmt.upper().startswith("CREATE DATABASE"):
                            statements.append(stmt)
                    
                    for stmt in statements:
                        cur.execute(stmt)
                    conn.commit()
                    logger.debug("Database schema synchronized successfully!")
                else:
                    logger.warning("scripts/init_db.sql not found at %s. Cannot synchronize schema.", schema_path)
            conn.close()
            logger.info("Database connection verified successfully!")
        except Exception as e:
            logger.error("Database connection/initialization failed: %s", e)
    else:
        logger.warning("POSTGRES_URL environment variable is not set.")

    # 2. Verify Qdrant
    try:
        client = QdrantClient(url=qdrant_url, api_key=qdrant_api_key, check_compatibility=False)
        client.get_collections()
        logger.info("Qdrant connection verified successfully!")
    except Exception as e:
        logger.error("Qdrant connection failed: %s", e)

    # 3. Start Auto-Ingestion watch thread
    try:
        start_auto_ingestion_watcher()
    except Exception as e:
        logger.error("Failed to start auto-ingestion watcher: %s", e)

    # 4. Run Guardrail Self Test & Background Health Probe
    try:
        run_startup_self_test(_config)
        asyncio.create_task(background_health_loop(_config))
    except Exception as e:
        logger.critical("Startup checks or health loop start failed: %s", e)
        raise

# Loaded once at import; the registry caches model singletons across requests.
_config = load_config(CONFIG_PATH)
# Warm torch/nomic BEFORE anything can import paddle (paddleocr) — paddle-first
# corrupts torch's allocator. paddleocr is lazy-imported, so this ordering holds.
warm_up(_config)
# The OCR engine (surya|paddle, config `vision_ocr`/scanned-path settings) is NOT
# warmed here — only the warm_up() ordering above is what prevents the
# paddle/torch allocator crash. (An earlier version of this comment described
# running OCR in an isolated subprocess; that isolation does not exist in the
# current code — verify the crash risk is still mitigated before changing
# warm_up() ordering or the OCR call path.)
_registry = build_default_registry()
_agent_registry = build_agent_registry()
# In-memory per-session LangChain message history for the agent chat — good enough
# for a demo; resets on restart. Only advanced past a turn that fully completed
# (not one awaiting write approval), so a decline/retry replays cleanly.
_agent_sessions: dict[str, list] = {}

# search_documents/ingest_document results can be large (full citation snippets,
# table_data, image paths...), and every intermediate tool-call round-trip adds
# more messages on top. tools_node needs the FULL detail for the current turn
# (it's what the frontend's citation strip parses from tool_calls[].result), but
# once the turn is done we only cache the plain Q&A — not the tool calls/results
# or the system prompt — for replay as conversation_history on the NEXT question.
# The model only needs to remember what was asked and answered, not how it got
# there; re-sending full tool payloads on every later turn is exactly what was
# blowing through Groq free-tier TPM/TPD limits after just one or two follow-ups.
# Same simplification _history_to_messages already applies to history reloaded
# from Postgres — this just applies it to the in-memory cache too.
def _qa_only(messages: list) -> list:
    """Filter messages to retain strictly clean User Questions (HumanMessage)
    and final Assistant Answers (AIMessage without tool calls). Excludes intermediate
    tool messages, system prompts, and raw execution logs.
    """
    clean = []
    for m in messages:
        if isinstance(m, HumanMessage):
            text = clean_message_content(m.content).strip()
            if text:
                clean.append(HumanMessage(content=text))
        elif isinstance(m, AIMessage) and not getattr(m, "tool_calls", None):
            text = clean_message_content(m.content).strip()
            if text:
                clean.append(AIMessage(content=text))
    return clean




CONFIG_DIR = os.path.dirname(CONFIG_PATH) or "config"


def _reload_pipeline() -> None:
    """Re-read CONFIG_PATH and rebuild the pipeline objects in place after a config
    edit. Models are NOT re-warmed (routing/prompt/OCR/chunking edits don't change
    the embedding/vision models); to switch those, edit + restart the server.
    ingest_document derives the ingestion profile from _config on each call, so the
    reloaded _config is all it needs — nothing else to rebuild here."""
    global _config, _registry
    _config = load_config(CONFIG_PATH)
    # OCR engine/timeout are read from config by the isolated OCR subprocess at
    # ingest time, so there's nothing to set in-process here.
    _registry = build_default_registry()


class AgentChatRequest(BaseModel):
    message: str
    session_id: str = "web"
    message_id: str = ""  # Client-generated UUID for deduplication
    approved_writes: bool = False
    # the pending write(s) the user approved, echoed back so approval is bound to the
    # exact name+args that were shown (not whatever the model re-proposes).
    approved_calls: list[dict] | None = None
    active_document_id: str | None = None  # Active document open in frontend viewer



class ConfigSave(BaseModel):
    yaml: str


class ProfileSave(BaseModel):
    name: str
    yaml: str


class ProfileActivate(BaseModel):
    name: str


class SettingsSave(BaseModel):
    settings: dict
    save_as: str | None = None   # optional: write to a new profile instead of active


# Flat form-field -> nested config path. The Settings UI edits these; everything
# else stays as-is in the YAML.
_SETTINGS_MAP = {
    "default_industry": ["default_industry"],
    "pdf_extractor_digital": ["pdf_extractors", "digital"],
    "ocr_engine": ["ocr", "engine"],
    # global LLM defaults
    "llm_provider": ["llm", "provider"],
    "llm_model": ["llm", "model"],
    "llm_answer_model": ["llm", "answer_model"],
    # vision
    "vision_provider": ["vision", "provider"],
    "vision_model": ["vision", "model"],
    "vision_ocr_model": ["vision_ocr", "model"],
    "vision_enabled": ["vision", "enabled"],
    # agent
    "agent_provider": ["query", "agent", "provider"],
    "agent_model": ["query", "agent", "model"],
    # per-step model overrides (blank => inherit global llm.model)
    "categorization_model": ["categorization", "model"],
    "enrichment_model": ["enrichment", "model"],
    "planner_model": ["query", "planner", "model"],
    "answerer_model": ["query", "answerer", "model"],
    # chunking / enrichment
    "chunking_strategy": ["chunking", "strategy"],
    "chunking_size": ["chunking", "size"],
    "chunking_overlap": ["chunking", "overlap"],
    "enrichment_summarize": ["enrichment", "summarize"],
    "enrichment_keyword_count": ["enrichment", "keyword_count"],
    "enrichment_prompt": ["enrichment", "prompt"],
    # embeddings & reranking
    "embeddings_dense_provider": ["embeddings", "dense_provider"],
    "embeddings_dense_model": ["embeddings", "dense_model"],
    "embeddings_dense_dim": ["embeddings", "dense_dim"],
    "embeddings_reranker_provider": ["embeddings", "reranker_provider"],
    "embeddings_reranker_model": ["embeddings", "reranker_model"],
    # auto-ingestion
    "auto_ingestion_enabled": ["auto_ingestion", "enabled"],
    "auto_ingestion_watch_dir": ["auto_ingestion", "watch_dir"],
    "auto_ingestion_poll_interval": ["auto_ingestion", "poll_interval"],
    "auto_ingestion_on_success": ["auto_ingestion", "on_success"],
    "auto_ingestion_on_failure": ["auto_ingestion", "on_failure"],
    # docling extraction mode
    "docling_mode": ["extraction", "docling", "mode"],
    "docling_server_url": ["extraction", "docling", "server_url"],
    "docling_server_key": ["extraction", "docling", "server_key"],
    # storage provider settings
    "storage_provider": ["storage", "provider"],
    "supabase_url": ["storage", "supabase_url"],
    "supabase_key": ["storage", "supabase_key"],
    "supabase_bucket": ["storage", "supabase_bucket"],
}

# Optional model-override fields: a BLANK value means "inherit the global llm block",
# so on save we REMOVE the key rather than writing an empty model name.
_OPTIONAL_OVERRIDE_KEYS = {
    "llm_answer_model", "vision_ocr_model", "agent_provider", "agent_model",
    "categorization_model", "enrichment_model", "planner_model", "answerer_model",
}


def _settings_view(cfg: dict) -> dict:
    """Curated, flat view of the editable settings for the form UI."""
    def dig(path, default=None):
        cur = cfg
        for k in path:
            if not isinstance(cur, dict):
                return default
            cur = cur.get(k)
        return cur if cur is not None else default
    load_dotenv(override=True)
    view = {k: dig(p) for k, p in _SETTINGS_MAP.items()}
    # Resolve env var placeholder in docling_server_url if present
    durl = str(view.get("docling_server_url") or "").strip()
    if durl.startswith("${") and durl.endswith("}"):
        var_name = durl[2:-1].strip()
        if var_name.startswith("http://") or var_name.startswith("https://"):
            view["docling_server_url"] = var_name
        else:
            view["docling_server_url"] = os.getenv(var_name) or os.getenv("DOCLING_SERVER_URL", "http://localhost:8083")
    elif not durl:
        view["docling_server_url"] = os.getenv("DOCLING_SERVER_URL", "http://localhost:8083")

    # structured (dict/list) settings edited with dedicated UI widgets
    view["vision_prompts"] = dig(["vision", "prompt"]) or {}
    view["ingestion_steps"] = dig(["ingestion", "steps"]) or []
    view["route_gates"] = dig(["ingestion", "route_gates"]) or {}
    # dropdown / option lists
    view["_industries"] = cfg.get("industries", [])
    view["_available_tools"] = sorted(_registry.names())
    view["_digital_pdf_options"] = ["pymupdf_pdf", "docling_pdf"]
    view["_active"] = os.path.basename(CONFIG_PATH)
    return view


def _apply_settings(raw: dict, settings: dict) -> dict:
    return apply_settings_in_place(raw, settings, _SETTINGS_MAP, _OPTIONAL_OVERRIDE_KEYS)


def _validate_yaml(text: str) -> dict:
    try:
        parsed = yaml.safe_load(text)
    except Exception as exc:
        raise HTTPException(400, f"invalid YAML: {exc}") from exc
    if not isinstance(parsed, dict):
        raise HTTPException(400, "config must be a YAML mapping")
    return parsed


def _file_type(filename: str) -> str:
    return EXT_TO_FILE_TYPE.get(os.path.splitext(filename)[1].lower(), "unknown")


# ── Postgres connection pool ──────────────────────────────────────────────────
# Reuses connections across requests instead of opening a new TCP connection
# per poll. Eliminates the dominant source of Supabase PostgREST log spam
# (previously: ~137 new connections/min from polling alone).
import threading as _threading
import queue as _queue
import psycopg as _psycopg

_PG_POOL_MIN = 2
_PG_POOL_MAX = 10
_pg_pool_lock = _threading.Lock()
_pg_pool_queue: "_queue.Queue[_psycopg.Connection]" = _queue.Queue(maxsize=_PG_POOL_MAX)
_pg_pool_count = 0  # total connections created


def _pool_get_conn() -> "_psycopg.Connection":
    """Get a healthy connection from the pool, creating one if below max."""
    global _pg_pool_count
    # Try to get an existing connection (non-blocking)
    try:
        conn = _pg_pool_queue.get_nowait()
        # Test it's still alive
        try:
            conn.execute("SELECT 1")
            return conn
        except Exception:
            try:
                conn.close()
            except Exception:
                pass
            with _pg_pool_lock:
                _pg_pool_count -= 1
    except _queue.Empty:
        pass

    # Create a new connection if under max
    with _pg_pool_lock:
        if _pg_pool_count < _PG_POOL_MAX:
            _pg_pool_count += 1
        else:
            raise RuntimeError("Connection pool exhausted")

    conn = _psycopg.connect(_pg_dsn(), autocommit=True, prepare_threshold=None)
    return conn


def _pool_return_conn(conn: "_psycopg.Connection") -> None:
    """Return a connection to the pool, or close it if pool is full."""
    global _pg_pool_count
    try:
        if not conn.closed:
            _pg_pool_queue.put_nowait(conn)
            return
    except (_queue.Full, Exception):
        pass
    # Pool full or connection dead — close it
    try:
        conn.close()
    except Exception:
        pass
    with _pg_pool_lock:
        _pg_pool_count -= 1


class _PooledPostgresStore(PostgresStore):
    """PostgresStore backed by the pool. close() returns the connection to the
    pool instead of closing it — safe to call in a finally block as usual."""

    def __init__(self) -> None:  # type: ignore[override]
        self.conn = _pool_get_conn()
        # Migration flag is already class-level; no DDL needed after first run

    def close(self) -> None:  # type: ignore[override]
        _pool_return_conn(self.conn)


def _pg() -> PostgresStore:
    try:
        return _PooledPostgresStore()
    except Exception as exc:
        raise HTTPException(503, f"database unavailable: {exc}") from exc



# Ingestion runs in a FastAPI background task so /upload returns immediately. Progress
# is written to Postgres after EVERY step — the documents row is the single source of
# truth — so /files/{id}/progress reads it straight from the DB. No in-memory state:
# survives restarts and works across multiple workers.
_INGEST_STEPS = list((_config.get("ingestion") or {}).get("steps") or [])


def _save_page_images(document_id: str, pdf_path: str, pages: set, dpi: int = 150) -> list:
    """Render the given (1-based) PDF pages to JPEGs under uploads/pages/<doc_id>/
    and return [(page, web_path, width, height)]. PDF-only; best-effort. fitz
    rendering is pure C (no torch/paddle), so it's safe in the backend process."""
    import fitz
    page_dir = os.path.join(_PAGES_DIR, document_id)
    os.makedirs(page_dir, exist_ok=True)
    out = []
    doc = fitz.open(pdf_path)
    try:
        for p in sorted(pages):
            if p < 1 or p > len(doc):
                continue
            pix = doc[p - 1].get_pixmap(dpi=dpi)
            fname = f"p{p}.jpg"
            with open(os.path.join(page_dir, fname), "wb") as f:
                f.write(pix.tobytes("jpeg", jpg_quality=80))
            out.append((p, f"/pages/{document_id}/{fname}", pix.width, pix.height))
    finally:
        doc.close()
    return out


def _run_ingestion(document_id: str, dest: str, file_type: str, filename: str, session_id: str | None = None, message_id: str | None = None) -> None:
    """Run the full pipeline for one upload. Delegates run + status + finalize to the
    shared ingest_document entry point; the API-only tail (live DB progress + PDF
    page images) is injected via the on_step / on_complete hooks."""
    
    # Download staged/uploaded file if it starts with supabase://
    if dest.startswith("supabase://"):
        try:
            parts = dest[11:].split("/", 1)
            bucket = parts[0]
            key = parts[1]
            local_dest = os.path.join(UPLOAD_DIR, f"{document_id}_{filename}")
            from backend.storage.supabase_store import download_from_supabase
            download_from_supabase(bucket, key, local_dest, config=_config)
            dest = local_dest
        except Exception as e:
            logger.exception("Failed to download staged document from Supabase: %s", dest)
            if session_id and message_id:
                try:
                    from backend.storage.conversation_store import get_conversation_store
                    get_conversation_store().update_turn_by_message_id(
                        message_id,
                        content=f"Ingestion failed: Failed to download from Supabase storage: {str(e)}",
                        metadata={"type": "ingest_error", "filename": filename, "errorMsg": str(e), "message_id": message_id}
                    )
                except Exception:
                    pass
            raise e

    total = len(_INGEST_STEPS) or None

    # Pre-render slides if it is a PowerPoint file
    page_dir = os.path.join(_PAGES_DIR, document_id)
    if file_type == "ppt":
      from backend.core.office_renderer import render_pptx_slides
      try:
          render_pptx_slides(dest, page_dir)
      except Exception as e:
          logger.exception("Failed to render PPTX slides for %s", document_id)

    # One connection for the whole run (per-step progress UPDATEs + page images).
    pg = None
    try:
        pg = PostgresStore()
    except Exception:
        logger.exception("progress DB unavailable for %s; running without live progress",
                         document_id)

    def on_step(entry: dict, snapshot: dict) -> None:
        ms = entry.get("ms")
        dur = f"{ms / 1000:.1f}s" if isinstance(ms, (int, float)) else "?"
        logger.info("  · %-24s %-8s %s", entry.get('step', ''), entry.get('status', ''), dur)

        if pg is None:
            return
        metrics = snapshot.get("metrics", []) or []
        progress = min(len(metrics) / total, 0.99) if total else None
        try:
            pg.update_progress(
                document_id,
                metrics=metrics,
                current_step=entry.get("step"),
                progress=progress,
                total_steps=total,
                route=snapshot.get("route"),
                confidence=snapshot.get("confidence"),
                document_type=snapshot.get("document_type"),
                industry=snapshot.get("industry"),
            )
        except Exception:
            logger.warning("update_progress failed for %s (step=%s)",
                           document_id, entry.get("step"))

    def on_complete(result: dict) -> None:
        pass

    try:
        ingest_document(dest, document_id, config=_config, registry=_registry,
                        on_step=on_step, on_complete=on_complete)
        if session_id and message_id:
            try:
                from backend.storage.conversation_store import get_conversation_store
                get_conversation_store().update_turn_by_message_id(
                    message_id,
                    content=f"📎 {filename} ingested successfully!",
                    metadata={"type": "ingest_done", "filename": filename, "message_id": message_id}
                )
            except Exception:
                logger.exception("Failed to update direct ingest success status in DB")
    except Exception as e:
        if session_id and message_id:
            try:
                from backend.storage.conversation_store import get_conversation_store
                get_conversation_store().update_turn_by_message_id(
                    message_id,
                    content=f"Ingestion failed: {str(e)}",
                    metadata={"type": "ingest_error", "filename": filename, "errorMsg": str(e), "message_id": message_id}
                )
            except Exception:
                logger.exception("Failed to update direct ingest failure status in DB")
        raise e
    finally:
        if pg is not None:
            pg.close()


@app.post("/upload")
def upload(background: BackgroundTasks, file: UploadFile = File(...)):
    document_id = str(uuid.uuid4())
    os.makedirs(UPLOAD_DIR, exist_ok=True)
    dest = os.path.join(UPLOAD_DIR, f"{document_id}_{file.filename}")
    with open(dest, "wb") as f:
        shutil.copyfileobj(file.file, f)

    file_type = _file_type(file.filename)
    
    # Supabase upload check
    storage_cfg = _config.get("storage", {})
    provider = storage_cfg.get("provider", "local")
    if provider == "supabase":
        bucket = storage_cfg.get("supabase_bucket", "documents")
        key = f"{document_id}_{file.filename}"
        from backend.storage.supabase_store import upload_to_supabase
        try:
            db_path = upload_to_supabase(bucket, key, dest, config=_config)
        except Exception as e:
            logger.exception("Failed to upload to Supabase storage")
            raise HTTPException(500, f"Failed to upload to Supabase storage: {e}")
    else:
        db_path = dest

    pg = _pg()
    try:
        pg.insert_document(document_id, file.filename, file_type, db_path)
    finally:
        pg.close()

    # Kick off ingestion in the background; progress is persisted to Postgres and the
    # UI polls /files/{id}/progress (DB-backed) below.
    background.add_task(_run_ingestion, document_id, dest, file_type, file.filename)

    return {"id": document_id, "document_id": document_id,
            "filename": file.filename, "file_type": file_type,
            "status": "processing", "metrics": []}


@app.post("/files/stage")
def stage_file(file: UploadFile = File(...)):
    """Save an uploaded file to disk WITHOUT triggering ingestion — for the agent
    chat's file-attach flow. /upload always auto-ingests (deterministic path, no
    agent involved); a staged file instead becomes a file_path the agent can see
    and choose (with approval, since ingest_document is a write) to ingest via
    the normal agent-executor flow. Not registered as a `documents` row — it only
    becomes one if/when ingest_document actually runs on it."""
    os.makedirs(UPLOAD_DIR, exist_ok=True)
    stage_id = str(uuid.uuid4())
    dest = os.path.join(UPLOAD_DIR, f"{stage_id}_{file.filename}")
    with open(dest, "wb") as f:
        shutil.copyfileobj(file.file, f)

    storage_cfg = _config.get("storage", {})
    provider = storage_cfg.get("provider", "local")
    if provider == "supabase":
        bucket = storage_cfg.get("supabase_bucket", "documents")
        key = f"staged/{stage_id}_{file.filename}"
        from backend.storage.supabase_store import upload_to_supabase
        try:
            supabase_path = upload_to_supabase(bucket, key, dest, config=_config)
            return {"file_path": supabase_path, "filename": file.filename}
        except Exception as e:
            logger.exception("Failed to upload staged file to Supabase")
            raise HTTPException(500, f"Failed to upload staged file to Supabase: {e}")

    return {"file_path": dest, "filename": file.filename}


@app.post("/files/ingest-staged")
async def ingest_staged(body: dict, background: BackgroundTasks):
    """Trigger ingestion on an already-staged file (one previously saved by /files/stage).
    The chat UI calls this when the user drops a file with no question text — it bypasses
    the LLM agent entirely, so ingestion starts immediately after the user approves inline.
    The file bytes are already on disk; we just create the documents row and kick off the
    background pipeline."""
    file_path = body.get("file_path", "")
    filename = body.get("filename", "")
    session_id = body.get("session_id")
    message_id = body.get("message_id")

    is_supabase = file_path.startswith("supabase://")
    if not is_supabase and (not file_path or not os.path.isfile(file_path)):
        raise HTTPException(404, "staged file not found on disk")
    if not filename:
        filename = os.path.basename(file_path)
        # Strip the stage_id prefix (format: <uuid>_<realname>) if present
        parts = filename.split("_", 1)
        if len(parts) == 2 and len(parts[0]) == 36:  # UUID length
            filename = parts[1]

    document_id = str(uuid.uuid4())
    file_type = _file_type(filename)

    pg = _pg()
    try:
        pg.insert_document(document_id, filename, file_type, file_path)
    finally:
        pg.close()

    if session_id and message_id:
        try:
            from backend.storage.conversation_store import get_conversation_store
            get_conversation_store().update_turn_by_message_id(
                message_id,
                content="",
                metadata={
                    "type": "ingest_progress",
                    "filename": filename,
                    "stagedPath": file_path,
                    "documentId": document_id,
                    "message_id": message_id,
                }
            )
        except Exception:
            logger.exception("Failed to update direct ingest progress status in DB")

    background.add_task(_run_ingestion, document_id, file_path, file_type, filename, session_id, message_id)

    return {
        "id": document_id,
        "document_id": document_id,
        "filename": filename,
        "file_type": file_type,
        "status": "processing",
        "metrics": [],
    }


@app.post("/agent/sessions/{session_id}/init-direct-ingest")
def init_direct_ingest(session_id: str, body: dict):
    filename = body.get("filename")
    staged_path = body.get("staged_path")
    message_id = body.get("message_id")
    if not message_id:
        message_id = str(uuid.uuid4())
    
    from backend.storage.conversation_store import get_conversation_store
    store = get_conversation_store()
    
    # Save user turn
    store.save_turn(session_id, "user", f"📎 {filename}")
    
    # Save assistant turn (the ingest approval card)
    store.save_turn(
        session_id,
        "assistant",
        "",
        {"type": "ingest_approval", "filename": filename, "stagedPath": staged_path, "message_id": message_id}
    )
    return {"message_id": message_id}


@app.post("/agent/sessions/{session_id}/cancel-direct-ingest")
def cancel_direct_ingest(session_id: str, body: dict):
    message_id = body.get("message_id")
    filename = body.get("filename")
    if message_id:
        from backend.storage.conversation_store import get_conversation_store
        get_conversation_store().update_turn_by_message_id(
            message_id,
            content="Ingestion cancelled.",
            metadata={"type": "ingest_cancelled", "filename": filename, "message_id": message_id}
        )
    return {"ok": True}


@app.get("/files/{file_id}/progress")
def file_progress(file_id: str):
    """Live per-step status for an in-flight (or finished) ingestion, read straight
    from Postgres — the single source of truth (survives restarts + multi-worker)."""
    pg = _pg()
    try:
        doc = pg.get_document(file_id)
    finally:
        pg.close()
    if not doc:
        raise HTTPException(404, "document not found")
    return doc


@app.get("/files")
def list_files():
    pg = _pg()
    try:
        return pg.list_documents()
    finally:
        pg.close()


@app.get("/files/{file_id}")
def get_file(file_id: str):
    pg = _pg()
    try:
        doc = pg.get_document(file_id)
    finally:
        pg.close()
    if not doc:
        raise HTTPException(404, "document not found")
    return doc


@app.get("/files/{file_id}/pages")
def file_pages(file_id: str):
    """List rendered full-page images for a document (page -> /pages/... url).
    Only pages that produced chunks are present."""
    pg = _pg()
    try:
        return {"document_id": file_id, "pages": pg.list_page_images(file_id)}
    finally:
        pg.close()


@app.get("/files/{file_id}/pages/{page}")
def file_page(file_id: str, page: int):
    """Metadata for one page image (bytes served at the returned image_path).
    Lets the answerer/agent pull up a full page for visual grounding."""
    pg = _pg()
    try:
        rec = pg.get_page_image(file_id, page)
    finally:
        pg.close()
    if not rec:
        raise HTTPException(404, "page image not found")
    return rec


def _ensure_local_file(document_id: str, doc: dict) -> str:
    """Helper to ensure a document's original file is present on local disk.
    If the file is backed by Supabase storage (stored as supabase:// URI in file_path)
    and is missing from local disk, it is downloaded from Supabase to UPLOAD_DIR.
    
    Returns the resolved local file path.
    """
    file_path = doc.get("file_path") or ""
    filename = doc.get("filename") or "document"
    
    if file_path.startswith("supabase://"):
        local_dest = os.path.join(UPLOAD_DIR, f"{document_id}_{filename}")
        if not os.path.isfile(local_dest):
            try:
                parts = file_path[11:].split("/", 1)
                bucket = parts[0]
                key = parts[1]
                from backend.storage.supabase_store import download_from_supabase
                download_from_supabase(bucket, key, local_dest, config=_config)
            except Exception as e:
                logger.exception("Failed to download file from Supabase: %s", file_path)
                raise HTTPException(500, f"Failed to download file from Supabase storage: {e}")
        return local_dest
        
    if not file_path or not os.path.isfile(file_path):
        import glob as _glob
        matches = _glob.glob(os.path.join(UPLOAD_DIR, f"{document_id}_*"))
        if not matches:
            raise HTTPException(404, "Original document file not found")
        file_path = matches[0]
        
    return file_path


@app.get("/files/{file_id}/pdf")
def file_pdf(file_id: str):
    """Stream the original PDF file for this document so the frontend can
    embed it in an <iframe> with a #page=N fragment for visual grounding.
    Only serves PDF documents; returns 404 for other file types."""
    from fastapi.responses import FileResponse
    pg = _pg()
    try:
        doc = pg.get_document(file_id)
    finally:
        pg.close()
    if not doc:
        raise HTTPException(404, "document not found")
    if doc.get("file_type") != "pdf":
        raise HTTPException(400, "document is not a PDF")
    
    file_path = _ensure_local_file(file_id, doc)
    filename = doc.get("filename") or os.path.basename(file_path)
    return FileResponse(
        path=file_path,
        media_type="application/pdf",
        filename=filename,
        headers={"Content-Disposition": f'inline; filename="{filename}"'},
    )


@app.get("/files/{file_id}/raw")
def file_raw(file_id: str):
    """Serve the raw file for this document (e.g. PDF, XLSX, PPTX) so the client can download or parse it."""
    from fastapi.responses import FileResponse
    import mimetypes
    pg = _pg()
    try:
        doc = pg.get_document(file_id)
    finally:
        pg.close()
    if not doc:
        raise HTTPException(404, "document not found")
    
    file_path = _ensure_local_file(file_id, doc)
    filename = doc.get("filename") or os.path.basename(file_path)
    media_type, _ = mimetypes.guess_type(file_path)
    media_type = media_type or "application/octet-stream"

    return FileResponse(
        path=file_path,
        media_type=media_type,
        filename=filename,
        headers={"Content-Disposition": f'inline; filename="{filename}"'},
    )


@app.get("/files/{file_id}/pages/{page}/image")
def file_page_image_ondemand(file_id: str, page: int):
    """Render a specific page of a PDF or serve an exported PPT slide image."""
    from fastapi.responses import FileResponse, Response
    import fitz

    pg = _pg()
    try:
        doc = pg.get_document(file_id)
    finally:
        pg.close()

    if not doc:
        raise HTTPException(404, "document not found")

    file_type = doc.get("file_type") or ""
    file_path = _ensure_local_file(file_id, doc)

    if file_type == "ppt":
        page_dir = os.path.join(_PAGES_DIR, file_id)
        img_path = os.path.join(page_dir, f"p{page}.jpg")
        if not os.path.isfile(img_path):
            from backend.core.office_renderer import render_pptx_slides
            logger.info("On-demand slide rendering triggered for PPT: %s", file_path)
            try:
                render_pptx_slides(file_path, page_dir)
            except Exception as render_err:
                logger.exception("Failed on-demand PPT rendering: %s", render_err)
        if os.path.isfile(img_path):
            return FileResponse(img_path, media_type="image/jpeg")
        raise HTTPException(404, f"Slide image not found: page {page}")

    try:
        pdf = fitz.open(file_path)
        total_pages = len(pdf)
        if page < 1 or page > total_pages:
            raise HTTPException(400, f"invalid page number: {page}. PDF has {total_pages} pages.")

        # Render page to JPEG with 150 DPI for high quality
        pix = pdf[page - 1].get_pixmap(dpi=150)
        image_bytes = pix.tobytes("jpeg", jpg_quality=85)
        pdf.close()

        return Response(content=image_bytes, media_type="image/jpeg")
    except HTTPException:
        raise
    except Exception as exc:
        logger.exception("Failed to render page %d on-demand for %s", page, file_id)
        raise HTTPException(500, f"failed to render page: {exc}")


@app.get("/files/{file_id}/pdf-info")
def file_pdf_info(file_id: str):
    """Get metadata about the PDF (like total number of pages) using PyMuPDF or slide count for PPT."""
    import fitz
    pg = _pg()
    try:
        doc = pg.get_document(file_id)
    finally:
        pg.close()
    if not doc:
        raise HTTPException(404, "document not found")

    file_type = doc.get("file_type") or ""
    file_path = _ensure_local_file(file_id, doc)

    if file_type == "ppt":
        page_dir = os.path.join(_PAGES_DIR, file_id)
        if os.path.isdir(page_dir):
            import glob
            files = glob.glob(os.path.join(page_dir, "p*.jpg"))
            if files:
                return {"total_pages": len(files)}
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            return {"total_pages": len(prs.slides)}
        except Exception:
            pass
        return {"total_pages": 0}

    try:
        pdf = fitz.open(file_path)
        total_pages = len(pdf)
        pdf.close()
        return {"total_pages": total_pages}
    except Exception as exc:
        logger.exception("Failed to get PDF info for %s", file_id)
        raise HTTPException(500, f"failed to read PDF: {exc}")


@app.get("/files/{file_id}/original")
def file_original(file_id: str):
    """Stream the original file bytes for types that need no transformation to
    view — currently just images. Browsers render image bytes natively, unlike
    PDF (needs page rasterization, see file_pdf/file_page_image_ondemand above)
    or docx (needs HTML conversion, see file_docx_html below)."""
    from fastapi.responses import FileResponse
    import mimetypes
    pg = _pg()
    try:
        doc = pg.get_document(file_id)
    finally:
        pg.close()
    if not doc:
        raise HTTPException(404, "document not found")
    if doc.get("file_type") != "image":
        raise HTTPException(400, "document is not an image")

    file_path = _ensure_local_file(file_id, doc)
    filename = doc.get("filename") or os.path.basename(file_path)
    media_type = mimetypes.guess_type(filename)[0] or "application/octet-stream"
    return FileResponse(
        path=file_path,
        media_type=media_type,
        filename=filename,
        headers={"Content-Disposition": f'inline; filename="{filename}"'},
    )


@app.get("/files/{file_id}/docx-html")
def file_docx_html(file_id: str):
    """Convert a docx to HTML for in-panel viewing.

    Unlike PDF, Word has no fixed, ingest-time-knowable page number — pagination
    is renderer-dependent (fonts/margins/zoom all shift it). So there's no
    per-page endpoint here, and no page-image model: this returns the WHOLE
    document as HTML, and the frontend (PageViewerPanel) matches a citation to
    a location in it at VIEW TIME by searching for the citation's own snippet
    text, rather than jumping to a persisted page/paragraph number. See
    backend/extraction/word/tool.py for why paragraph_index is deliberately
    NOT threaded into source_ref/tags for this — the team's existing
    convention (see schemas.py's Excel cell_range note) keeps format-specific
    locators out of the shared citation schema.
    """
    pg = _pg()
    try:
        doc = pg.get_document(file_id)
    finally:
        pg.close()
    if not doc:
        raise HTTPException(404, "document not found")
    if doc.get("file_type") != "docx":
        raise HTTPException(400, "document is not a docx")

    file_path = _ensure_local_file(file_id, doc)

    try:
        import mammoth
        with open(file_path, "rb") as f:
            result = mammoth.convert_to_html(f)
    except Exception as exc:
        logger.exception("docx->html conversion failed for %s", file_id)
        raise HTTPException(500, f"failed to convert docx to HTML: {exc}")

    return {
        "html": result.value,
        "warnings": [str(w) for w in (result.messages or [])],
        "filename": doc.get("filename"),
    }


@app.delete("/files/{file_id}")
def delete_file(file_id: str):
    dim = _config.get("embeddings", {}).get("dense_dim", 768)
    collection = _config.get("database", {}).get("qdrant_collection", "chunks")

    vectors = QdrantStore(dim, collection)
    try:
        vectors.delete_by_document(file_id)
    finally:
        vectors.close()

    pg = _pg()
    file_path = None
    try:
        doc = pg.get_document(file_id)
        if doc:
            file_path = doc.get("file_path")
        pg.delete_document(file_id)
    finally:
        pg.close()

    # Delete from Supabase storage if file is stored there
    if file_path and file_path.startswith("supabase://"):
        try:
            parts = file_path[11:].split("/", 1)
            bucket = parts[0]
            key = parts[1]
            from backend.storage.supabase_store import delete_from_supabase
            delete_from_supabase(bucket, key, config=_config)
        except Exception:
            logger.exception("Failed to delete document from Supabase storage: %s", file_path)

    return {"deleted": file_id}


def _history_to_messages(history: list[dict]) -> list:
    """DB rows (role/content) -> LangChain messages, for seeding run_agent's
    conversation_history when a session isn't in the in-memory cache (server
    restart, or reopening an old chat). Drops tool-call structure — the model
    only needs the visible text to keep context, not the exact prior calls."""
    msgs = []
    for h in history:
        if h.get("role") == "user":
            msgs.append(HumanMessage(h.get("content") or ""))
        elif h.get("role") == "assistant" and h.get("content"):
            msgs.append(AIMessage(h["content"]))
    return msgs




@app.post("/agent/chat")
def agent_chat(req: AgentChatRequest, response: Response):
    """Agentic chat: the model picks which tool to call (ingest/search/sql) instead
    of always going straight to retrieval. Writes (ingest_document) stop and report
    what they want to run — POST again with approved_writes=true to actually run it.
    Turns are persisted (Postgres) so the UI's session sidebar can list + reopen
    past conversations; the in-memory cache is just a fast path within one run.
    """
    from backend.storage.conversation_store import get_conversation_store

    # 1. Token Quota Check (Reserve budget)
    quota = get_enforcer(_config)
    reserve = get_reserve_tokens(_config)
    if not quota.reserve(req.session_id, reserve):
        raise HTTPException(429, "Token budget exceeded for this session. Please try again later.")

    max_history = (_config.get("query", {}).get("agent", {}) or {}).get("max_history_messages", 20)
    history = _agent_sessions.get(req.session_id)
    if history is None:
        try:
            history = _history_to_messages(
                get_conversation_store().load_history(req.session_id, n=max_history * 2)
            )
        except Exception:
            logger.debug("agent chat history load failed", exc_info=True)
            history = []

    # Ensure clean Q&A pairs only and slice to max_history (strictly last 10 Q&A pairs / 20 messages max)
    history = _qa_only(history)[-max_history:]

    # Save user message immediately so it's persisted and visible if user switches/reloads.
    # The [Attached file: <name> — path: <disk_path>] annotation is an internal protocol
    # marker sent to the agent; replace it with "📎 <name>" so reloaded history shows just
    # the icon + filename. Also strip the legacy "Please ingest this file." sentinel that
    # old builds inserted when the user typed nothing (the frontend no longer sends it).
    import re as _re
    def _to_display(m):
        fname = m.group(1).strip()
        return f"\n\n📎 {fname}"
    _display_message = _re.sub(
        r"\s*\[Attached file:\s*(.+?)\s*(?:—|-)\s*path:[^\]]*\]",
        _to_display,
        req.message,
    )
    # Strip the now-obsolete sentinel (may appear in messages from older clients)
    _display_message = _re.sub(r"^Please ingest this file\.\s*", "", _display_message, flags=_re.IGNORECASE)
    _display_message = _display_message.strip() or req.message
    try:
        store = get_conversation_store()
        # Guard against duplicate user messages: if the most recent turn in this session
        # is already a user message with the same text, don't save again.  This prevents
        # the double-bubble when the user navigates away and the ChatPage reloads, which
        # can cause the same /agent/chat to be re-submitted.
        recent = store.load_history(req.session_id, n=1)
        if not (recent and recent[-1].get("role") == "user" and recent[-1].get("content") == _display_message):
            store.save_turn(req.session_id, "user", _display_message)
    except Exception:
        logger.debug("agent user chat history save failed", exc_info=True)

    result = run_agent(
        req.message, config=_config, registry=_agent_registry,
        conversation_history=history, approved_writes=req.approved_writes,
        approved_calls=req.approved_calls, session_id=req.session_id,
        active_document_id=req.active_document_id,
    )

    # 2. Reconcile token budget
    actual_tokens = result.get("token_usage", {}).get("total", 0)
    quota.reconcile(req.session_id, reserved=reserve, actual=actual_tokens)

    # 3. Add guardrail headers
    response.headers["X-Guardrail-Events"] = str(result.get("guard_policy", "allow"))
    response.headers["X-Risk-Score"] = str(result.get("guard_risk_score", 0))

    tool_calls = [
        {"name": c["name"], "args": c["args"], "result": c.get("result")}
        for c in result.get("tool_calls", [])
    ]
    exec_trace = result.get("execution_trace") or []

    if result["status"] == "done":
        _agent_sessions[req.session_id] = _qa_only(result["messages"])[-max_history:]
        try:
            store = get_conversation_store()
            store.save_turn(
                req.session_id, "assistant", result.get("answer") or "",
                metadata={
                    "status": "done",
                    "tool_calls": tool_calls,
                    "llm_calls": result.get("llm_calls"),
                    "execution_trace": exec_trace,
                    "token_usage": result.get("token_usage"),
                    "trace_id": result.get("trace_id"),
                    "cad_diagrams": result.get("cad_diagrams"),
                },
            )
        except Exception:
            logger.debug("agent chat history save failed", exc_info=True)
    elif result["status"] == "needs_clarification":
        # Persist the question so the follow-up (the user's choice) carries context —
        # otherwise the agent gets a bare option token with no memory of what it asked.
        _agent_sessions[req.session_id] = _qa_only(result["messages"])[-max_history:]

        try:
            store = get_conversation_store()
            store.save_turn(
                req.session_id, "assistant",
                result.get("answer") or result.get("question") or "",
                metadata={
                    "status": "needs_clarification",
                    "question": result.get("question"),
                    "options": result.get("options") or [],
                    "tool_calls": tool_calls,
                    "llm_calls": result.get("llm_calls") or [],
                    "execution_trace": exec_trace,
                    "token_usage": result.get("token_usage"),
                    "trace_id": result.get("trace_id"),
                    "cad_diagrams": result.get("cad_diagrams"),
                },
            )
        except Exception:
            logger.debug("agent chat clarification save failed", exc_info=True)
    elif result["status"] == "needs_approval":
        try:
            store = get_conversation_store()
            store.save_turn(
                req.session_id, "assistant",
                result.get("answer") or "",
                metadata={
                    "status": "needs_approval",
                    "pending": result.get("pending") or [],
                    "tool_calls": tool_calls,
                    "llm_calls": result.get("llm_calls") or [],
                    "execution_trace": exec_trace,
                    "token_usage": result.get("token_usage"),
                    "trace_id": result.get("trace_id"),
                },
            )
        except Exception:
            logger.debug("agent chat approval save failed", exc_info=True)
    return {
        "status": result["status"],
        "answer": result.get("answer"),
        "pending": result.get("pending"),
        # needs_clarification: a machine-readable chooser for the UI
        "question": result.get("question"),
        "options": result.get("options"),
        "tool_calls": tool_calls,
        "llm_calls": result.get("llm_calls") or [],
        "execution_trace": exec_trace,
        "token_usage": result.get("token_usage"),
        "trace_id": result.get("trace_id"),
        "cad_diagrams": result.get("cad_diagrams"),
    }






@app.get("/agent/sessions")
def list_agent_sessions():
    """Past conversations for the chat UI's sidebar, most recently active first."""
    from backend.storage.conversation_store import get_conversation_store

    return get_conversation_store().list_sessions()


@app.get("/agent/sessions/{session_id}")
def get_agent_session(session_id: str):
    """One conversation's full turn history, in the same shape /agent/chat
    returns per turn, so the frontend can render live and reloaded messages
    with the same component."""
    from backend.storage.conversation_store import get_conversation_store

    history = get_conversation_store().load_history(session_id, n=500)
    out = []
    cached_options = None
    for h in history:
        meta = h.get("metadata") or {}
        item = {"role": h["role"], "content": h["content"], **meta}
        # If this turn was a clarification question and options are missing, restore them
        if item.get("role") == "assistant" and (
            item.get("status") == "needs_clarification" or
            item.get("content") == "Which document would you like to search?"
        ) and not item.get("options"):
            if cached_options is None:
                from backend.storage.postgres_store import PostgresStore
                _SEARCHABLE_EXTS = {".pdf", ".xlsx", ".xls", ".csv", ".docx", ".doc", ".txt", ".pptx", ".ppt"}
                pg = PostgresStore(config=_config)
                try:
                    docs = pg.list_documents()
                    seen = set()
                    cached_options = []
                    for d in docs:
                        fname = d.get("filename", "")
                        if not fname:
                            continue
                        ext = "." + fname.rsplit(".", 1)[-1].lower() if "." in fname else ""
                        if ext not in _SEARCHABLE_EXTS:
                            continue
                        if fname not in seen:
                            seen.add(fname)
                            cached_options.append(fname)
                finally:
                    pg.close()
            item["status"] = "needs_clarification"
            item["options"] = cached_options
            item["question"] = item.get("question") or item.get("content")
        out.append(item)
    return out


@app.patch("/agent/sessions/{session_id}")
def patch_agent_session(session_id: str, body: dict):
    """Update session metadata: { "title": "...", "pinned": true/false }."""
    from backend.storage.conversation_store import get_conversation_store

    get_conversation_store().update_session(
        session_id,
        title=body.get("title"),
        pinned=body.get("pinned"),
    )
    return {"updated": session_id}


@app.delete("/agent/sessions/{session_id}")
def delete_agent_session(session_id: str):
    from backend.storage.conversation_store import get_conversation_store

    get_conversation_store().delete_session(session_id)
    _agent_sessions.pop(session_id, None)
    return {"deleted": session_id}


@app.get("/health")
def health():
    return {"status": "ok", "tools": sorted(_registry.names())}


# ── Config / customer profiles (edited from the UI Settings page) ─────────────
# The raw YAML keeps ${ENV} placeholders for secrets, so it's safe to read/write
# over the API (real keys live only in .env and are resolved at load time).

@app.get("/config/raw")
def config_raw():
    """Active config file as editable YAML text (secrets stay as ${ENV} refs)."""
    with open(CONFIG_PATH, encoding="utf-8") as f:
        return {"path": CONFIG_PATH, "active": os.path.basename(CONFIG_PATH), "yaml": f.read()}


@app.put("/config")
def config_save(body: ConfigSave):
    """Validate + write the active config, then hot-reload the pipeline."""
    _validate_yaml(body.yaml)
    with open(CONFIG_PATH, "w", encoding="utf-8") as f:
        f.write(body.yaml)
    _reload_pipeline()
    return {"ok": True, "path": CONFIG_PATH}


@app.get("/config/profiles")
def config_profiles():
    """List available config profiles (config/*.yaml) + which one is active."""
    files = sorted(os.path.basename(p) for p in glob.glob(os.path.join(CONFIG_DIR, "*.yaml")))
    return {"active": os.path.basename(CONFIG_PATH), "profiles": files}


@app.post("/config/profiles")
def config_save_profile(body: ProfileSave):
    """Save the YAML as a new/!existing profile config/<name>.yaml (per customer)."""
    _validate_yaml(body.yaml)
    name = os.path.basename(body.name.strip())
    if not name:
        raise HTTPException(400, "profile name required")
    if not name.endswith(".yaml"):
        name += ".yaml"
    with open(os.path.join(CONFIG_DIR, name), "w", encoding="utf-8") as f:
        f.write(body.yaml)
    return {"ok": True, "name": name}


@app.get("/config/settings")
def config_settings():
    """Flat, form-friendly view of the editable settings (no YAML for the user)."""
    raw_config = load_yaml_roundtrip(CONFIG_PATH)
    return _settings_view(raw_config)


@app.put("/config/settings")
def config_settings_save(body: SettingsSave):
    """Apply form settings to the active config (or save_as a new profile) + reload.
    Preserves all comments, formatting, and non-mapped keys."""
    global CONFIG_PATH
    raw = load_yaml_roundtrip(CONFIG_PATH)
    raw = _apply_settings(raw, body.settings)

    target = CONFIG_PATH
    if body.save_as and body.save_as.strip():
        name = os.path.basename(body.save_as.strip())
        if not name.endswith(".yaml"):
            name += ".yaml"
        target = os.path.join(CONFIG_DIR, name)
    dump_yaml_roundtrip(raw, target)
    CONFIG_PATH = target
    _reload_pipeline()
    return {"ok": True, "active": os.path.basename(CONFIG_PATH)}


@app.post("/config/activate")
def config_activate(body: ProfileActivate):
    """Switch the active profile to config/<name>.yaml and hot-reload."""
    global CONFIG_PATH
    name = os.path.basename(body.name)
    path = os.path.join(CONFIG_DIR, name)
    if not os.path.exists(path):
        raise HTTPException(404, f"profile {name} not found")
    CONFIG_PATH = path
    _reload_pipeline()
    return {"ok": True, "active": name}


@app.get("/health/docling-server")
def health_docling_server(url: str | None = None, mode: str | None = None):
    """Ping the configured or requested remote Docling server and return its status."""
    load_dotenv(override=True)
    dcfg = (_config.get("extraction") or {}).get("docling") or {}
    selected_mode = (mode or dcfg.get("mode") or "local").lower()
    
    # If mode is not remote and no explicit URL was passed to test, report local mode
    if selected_mode != "remote" and not url:
        return {"mode": "local", "reachable": None, "message": "Local mode — no server to check"}

    server_url = str(url or dcfg.get("server_url") or os.environ.get("DOCLING_SERVER_URL", "")).strip()
    if server_url.startswith("${") and server_url.endswith("}"):
        inner = server_url[2:-1].strip()
        if inner.startswith("http://") or inner.startswith("https://"):
            server_url = inner
        else:
            server_url = os.environ.get(inner, os.environ.get("DOCLING_SERVER_URL", "http://localhost:8083"))
    elif "${" in server_url or not server_url:
        server_url = os.environ.get("DOCLING_SERVER_URL", "http://localhost:8083")

    if not server_url:
        return {"mode": selected_mode, "reachable": False, "message": "No server URL configured"}

    try:
        import requests as _req
        r = _req.get(f"{server_url.rstrip('/')}/health", timeout=5)
        r.raise_for_status()
        return {"mode": selected_mode, "reachable": True, "url": server_url, "detail": r.json()}
    except Exception as e:
        return {"mode": selected_mode, "reachable": False, "url": server_url, "message": str(e)}


@app.get("/llm/calls")
@app.get("/llm-calls")
def get_all_llm_calls(
    session_id: str | None = None,
    document_id: str | None = None,
    kind: str | None = None,
    limit: int = 100,
    offset: int = 0,
):
    """Retrieve all LLM and vision calls across the entire project (agent, answerer, query_planner, hyde, vision, categorize, enrichment, etc.)."""
    pg = _pg()
    try:
        return pg.get_llm_calls(
            document_id=document_id,
            session_id=session_id,
            kind=kind,
            limit=limit,
            offset=offset,
        )
    finally:
        pg.close()


@app.get("/agent/sessions/{session_id}/llm-calls")
def get_session_llm_calls(session_id: str, limit: int = 100):
    """Retrieve LLM calls associated with a specific agent chat session."""
    pg = _pg()
    try:
        return pg.get_llm_calls(session_id=session_id, limit=limit)
    finally:
        pg.close()


@app.get("/files/{file_id}/vision-calls")
def get_vision_calls(file_id: str, limit: int = 200):
    """Retrieve all LLM/vision calls made for this document."""
    pg = _pg()
    try:
        return pg.get_llm_calls(document_id=file_id, limit=limit)
    finally:
        pg.close()


# ── P0 Compliance & Observability Endpoints ──────────────────────────────────

@app.get("/audit/queries")
def list_query_audits(
    request: Request,
    session_id: str | None = None,
    limit: int = 50,
    offset: int = 0,
):
    """Compliance: list query audit records (query text, retrieved chunk IDs, latency).

    Accessible to admin users. Pass X-Admin-Key header matching the ADMIN_API_KEY
    env var. Returns up to 50 records. Filter by session_id for a single conversation.
    """
    admin_key = os.getenv("ADMIN_API_KEY")
    if admin_key:
        provided = request.headers.get("X-Admin-Key", "")
        if provided != admin_key:
            raise HTTPException(status_code=403, detail="Admin access required")

    pg = _pg()
    try:
        return pg.list_query_audits(session_id=session_id, limit=limit, offset=offset)
    finally:
        pg.close()


@app.get("/audit/index-versions")
def list_index_versions(request: Request):
    """List registered index versions (embedding model + config hash + date).

    Lets you trace which embedding model was active when an answer was generated.
    Admin-only.
    """
    admin_key = os.getenv("ADMIN_API_KEY")
    if admin_key:
        provided = request.headers.get("X-Admin-Key", "")
        if provided != admin_key:
            raise HTTPException(status_code=403, detail="Admin access required")

    pg = _pg()
    try:
        rows = pg.conn.execute(
            "SELECT index_version, model_name, config_hash, created_at "
            "FROM index_versions ORDER BY created_at DESC"
        ).fetchall()
        return [
            {
                "index_version": r[0],
                "model_name": r[1],
                "config_hash": r[2],
                "created_at": r[3].isoformat() if r[3] else None,
            }
            for r in rows
        ]
    finally:
        pg.close()


@app.get("/")
def root():
    # Reload trigger comment v4
    return {"service": "Document Intelligence + RAG Accelerator", "docs": "/docs"}