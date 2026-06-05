"""Text extraction helpers for multiple file types (PDF, Excel, PPT, images).

Used for:
- industry keyword scoring (first 3 pages)
- TOC extraction (no API cost)
- CAD document detection and extraction
- Specialized engineering metadata extraction

The categorization engine still relies on vision for document_type.
"""

from __future__ import annotations

from typing import List, Dict, Any, Optional
import os
import re

import fitz

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from PIL import Image
    import pytesseract
    HAS_PYTESSERACT = True
except ImportError:
    HAS_PYTESSERACT = False


def extract_text(file_path: str, max_pages: int = 3) -> str:
    """Extract raw text from the first `max_pages` pages of any supported file type."""
    if not os.path.exists(file_path):
        return ""

    try:
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            return _extract_pdf_text(file_path, max_pages)
        elif ext in [".xlsx", ".xls"]:
            return _extract_excel_text(file_path, max_pages)
        elif ext in [".pptx", ".ppt"]:
            return _extract_ppt_text(file_path, max_pages)
        elif ext in [".png", ".jpg", ".jpeg", ".gif", ".bmp"]:
            return _extract_image_text(file_path)
        else:
            # For unknown types, try PDF first (might be misnamed)
            try:
                return _extract_pdf_text(file_path, max_pages)
            except Exception:
                return ""
    except Exception:
        return ""


def _extract_pdf_text(file_path: str, max_pages: int = 3) -> str:
    """Extract raw text from PDF pages."""
    text_parts: List[str] = []
    try:
        doc = fitz.open(file_path)
        pages_to_read = min(max_pages, len(doc))
        for page_num in range(pages_to_read):
            page = doc[page_num]
            text_parts.append(page.get_text())
        doc.close()
    except Exception:
        return ""
    return "\n".join(text_parts)


def _extract_excel_text(file_path: str, max_pages: int = 3) -> str:
    """Extract text from Excel workbook."""
    if not HAS_OPENPYXL:
        return ""

    text_parts: List[str] = []
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        sheet_count = 0
        for sheet in wb.sheetnames:
            if sheet_count >= max_pages:
                break
            ws = wb[sheet]
            text_parts.append(f"Sheet: {sheet}")
            for row in ws.iter_rows(values_only=True):
                text_parts.append(" | ".join(str(v) if v is not None else "" for v in row))
            sheet_count += 1
        wb.close()
    except Exception:
        pass
    return "\n".join(text_parts)


def _extract_ppt_text(file_path: str, max_pages: int = 3) -> str:
    """Extract text from PowerPoint presentation."""
    if not HAS_PPTX:
        return ""

    text_parts: List[str] = []
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides):
            if slide_num >= max_pages:
                break
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        text_parts.append(shape.text)
    except Exception:
        pass
    return "\n".join(text_parts)


def _extract_image_text(file_path: str) -> str:
    """Extract text from image using OCR (requires pytesseract)."""
    if not HAS_PYTESSERACT:
        return ""

    try:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        return text
    except Exception:
        return ""


def extract_toc_text(file_path: str, max_pages: int = 2) -> str:
    """Best-effort TOC extraction.

    If the PDF has a table of contents, fitz may expose it via doc.get_toc().
    We also do a minimal heuristic scan in the first few pages as a backup.
    """
    if not os.path.exists(file_path):
        return ""

    ext = os.path.splitext(file_path)[1].lower()

    # Only PDFs have TOC structure
    if ext != ".pdf":
        return extract_text(file_path, max_pages)

    try:
        doc = fitz.open(file_path)
        toc = []
        try:
            toc = doc.get_toc() or []
        except Exception:
            toc = []

        if toc:
            # toc entries: list of (level, title, page_number)
            lines = []
            for entry in toc:
                if len(entry) >= 2:
                    title = entry[1]
                    lines.append(str(title))
            doc.close()
            return "\n".join(lines)

        # fallback: scan first N pages for TOC-like patterns
        pages_to_read = min(max_pages, len(doc))
        parts: List[str] = []
        for i in range(pages_to_read):
            page = doc[i]
            parts.append(page.get_text())
        doc.close()
        return "\n".join(parts)

    except Exception:
        return ""


# ============================================================
# CAD DOCUMENT DETECTION & EXTRACTION
# ============================================================

CAD_KEYWORDS = [
    'drawing number', 'drawing no', 'drawing sheet', 'dwg',
    'model', 'scale', 'revision', 'rev', 'bom', 'bill of materials',
    'part number', 'part no', 'qty', 'quantity', 'description',
    'material', 'finish', 'tolerance', 'dimension', 'view',
    'section', 'detail', 'mm', 'inches', 'unless otherwise specified',
    'title block', 'signature block', 'approval', '図面番号',
    'scale 1:', 'drawing sheet', 'engineering drawing',
    'mechanical drawing', 'technical drawing', 'schematic',
    'cad', 'autocad', 'solidworks', 'catia', '3d model'
]

ENGINEERING_INDUSTRIES = {
    'automotive': ['vehicle', 'motor', 'engine', 'chassis', 'wiring', 'harness', 'ecu', 'drive', 'wheel', 'transmission'],
    'manufacturing': ['production', 'assembly', 'fixture', 'tolerance', 'qc', 'iso', 'jig', 'bore'],
    'engineering': ['voltage', 'resistor', 'schematic', 'pcb', 'torque', 'bearing', 'shaft', 'bolt', 'washer'],
    'aerospace': ['aircraft', 'avionics', 'flight', 'fuselage', 'landing', 'wing'],
    'pharma': ['equipment', 'validation', 'batch', 'sterile', 'controlled', 'chamber'],
}


def is_cad_document_by_filename(filename: str) -> bool:
    """
    Detect CAD documents by filename patterns.
    
    Looks for:
    - Part/drawing numbers: MS03AAA981AA, DWG-0001, 99Y_MKR2002100AB
    - Common CAD file naming: drawing_, schematic_, print_, layout_
    - Technical suffixes: -A, -B, -R1, -R2 (revisions)
    """
    lower = filename.lower()
    
    # Part number patterns (aerospace, automotive, manufacturing)
    part_number_patterns = [
        r'^[A-Z]{1,3}\d{2}[A-Z]{3}\d{6}',  # MS03AAA981AA style
        r'dwg[-_]?\d{4}',  # DWG-0001
        r'\d{2}y[-_]?[a-z]{3}\d{7}',  # 99Y_MKR2002100AB style
        r'[a-z]{2}\d{2}[a-z]{2}\d{3,6}[a-z]{2}',  # Generic part number
    ]
    
    for pattern in part_number_patterns:
        if re.search(pattern, lower):
            return True
    
    # Common CAD naming prefixes
    cad_prefixes = ['drawing_', 'schematic_', 'print_', 'layout_', 'cad_', 'dxf_', 'dwg_']
    if any(lower.startswith(p) for p in cad_prefixes):
        return True
    
    # Common CAD file suffixes and revisions
    cad_suffixes = ['-a.pdf', '-b.pdf', '-r1.pdf', '-r2.pdf', '-sheet.pdf', '_sheet.pdf']
    if any(lower.endswith(s) for s in cad_suffixes):
        return True
    
    # Motor/engine/mechanical keywords in filename
    mechanical_keywords = ['motor', 'engine', 'assembly', 'bracket', 'fixture', 'jig', 'bearing', 'shaft']
    word_count = sum(1 for kw in mechanical_keywords if kw in lower.replace('_', ' ').replace('-', ' '))
    if word_count >= 1:
        return True
    
    return False
    """Detect if PDF is engineering CAD drawing based on text keywords."""
    if not text or len(text) < 100:
        return False
    
    text_lower = text.lower()
    keyword_matches = sum(1 for kw in CAD_KEYWORDS if kw in text_lower)
    
    # At least 3 CAD keywords indicates CAD document
    return keyword_matches >= 3


def extract_cad_metadata(text: str) -> Dict[str, Optional[str]]:
    """Extract engineering metadata from CAD document text."""
    metadata = {
        'drawing_number': None,
        'title': None,
        'model': None,
        'scale': None,
        'industry': None,
    }
    
    text_lower = text.lower()
    
    # Extract drawing number
    drawing_patterns = [
        r'(?:drawing|drawing no|dwg)[\s:]*([A-Z0-9\-\/]+)',
        r'([A-Z]\d{2}[A-Z]\d{4}\d{3}[A-Z]{2})',  # Format like 99Y_MKR2002100AB
    ]
    
    for pattern in drawing_patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            metadata['drawing_number'] = match.group(1)
            break
    
    # Extract model
    model_match = re.search(r'(?:model|型式)[\s:]*([A-Z0-9\-]+)', text, re.IGNORECASE)
    if model_match:
        metadata['model'] = model_match.group(1)
    
    # Extract scale
    scale_match = re.search(r'scale[\s:]*([0-9:\.]+)', text, re.IGNORECASE)
    if scale_match:
        metadata['scale'] = scale_match.group(1)
    
    # Detect industry from keywords
    for industry, keywords in ENGINEERING_INDUSTRIES.items():
        if any(kw in text_lower for kw in keywords):
            metadata['industry'] = industry
            break
    
    # Fallback to 'engineering' if no specific industry detected
    if metadata['industry'] is None and is_cad_document(text):
        metadata['industry'] = 'engineering'
    
    return metadata


def extract_cad_bom(text: str) -> List[Dict[str, Any]]:
    """Extract BOM (Bill of Materials) from CAD document text."""
    bom_entries = []
    
    lines = text.split('\n')
    bom_section = False
    
    for line in lines:
        # Detect BOM section
        if any(kw in line.upper() for kw in ['BOM', 'BILL OF MATERIALS', '部品表']):
            bom_section = True
            continue
        
        if not bom_section:
            continue
        
        # Look for BOM entries: item number + part name + qty
        match = re.match(r'\s*(\d+)\s+(.+?)\s+(\d+)\s*(.*)', line)
        if match and len(match.group(2)) > 2:
            entry = {
                'item': match.group(1),
                'part_name': match.group(2).strip(),
                'qty': match.group(3),
                'notes': match.group(4).strip() if match.group(4) else '',
            }
            bom_entries.append(entry)
    
    return bom_entries


def extract_cad_dimensions(text: str) -> List[Dict[str, Any]]:
    """Extract dimensions from CAD document text."""
    dimensions = []
    
    # Pattern for dimensions: number + unit
    dim_pattern = r'(\d+(?:\.\d+)?)\s*(mm|cm|inch|in|″|\')'
    
    for match in re.finditer(dim_pattern, text):
        unit = match.group(2).replace('″', 'inch').replace('′', 'inch')
        dimensions.append({
            'value': float(match.group(1)),
            'unit': unit,
        })
    
    return dimensions


def detect_excel_document_type(file_path: str) -> tuple:
    """
    Detect the document type of an Excel file.
    
    Returns: (document_type, confidence, details)
    - document_type: invoice, financial_statement, purchase_order, budget, forecast, report (data analysis), spreadsheet
    - confidence: 0.0-1.0
    - details: dict with reasoning
    """
    if not HAS_OPENPYXL:
        return ("spreadsheet", 0.5, {"error": "openpyxl not available"})
    
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        all_text = []
        header_row = []
        sheet_names_lower = []
        
        # Scan all sheets for content patterns
        for sheet_idx, sheet_name in enumerate(wb.sheetnames[:3]):  # Check first 3 sheets
            ws = wb[sheet_name]
            sheet_names_lower.append(sheet_name.lower())
            
            # Get header row and content
            for idx, row in enumerate(ws.iter_rows(values_only=True)):
                if idx == 0:
                    header_row = [str(v).lower() if v else "" for v in row]
                all_text.append(" ".join(str(v).lower() if v else "" for v in row if v))
                if idx > 10:  # Sample first 10 rows
                    break
        
        wb.close()
        
        full_text = " ".join(all_text).lower()
        all_sheet_info = " ".join(sheet_names_lower + all_text).lower()
        
        # Excel document type keywords - enhanced
        doc_types = {
            "invoice": ["invoice", "inv-", "bill to", "amount due", "total amount", "payment", "customer"],
            "financial_statement": ["balance sheet", "income statement", "p&l", "profit", "loss", "assets", "liabilities", "equity", "fiscal"],
            "purchase_order": ["purchase order", "po ", "vendor", "supplier"],
            "budget": ["budget", "budgeted", "variance", "allocated"],
            "forecast": ["forecast", "projected", "projection", "estimate", "scenario"],
            "report": ["sales", "revenue", "quarterly", "monthly", "annual", "performance", "analysis"],
        }
        
        # Score each type
        scores = {}
        for doc_type, keywords in doc_types.items():
            matches = sum(1 for kw in keywords if kw in full_text or any(kw in h for h in header_row))
            scores[doc_type] = matches
        
        # Find best match with priority: prefer specific types over report
        best_type = None
        best_score = 0
        type_priority = ["invoice", "financial_statement", "purchase_order", "budget", "forecast", "report", "spreadsheet"]
        
        for doc_type in type_priority:
            if scores.get(doc_type, 0) > best_score:
                best_score = scores[doc_type]
                best_type = doc_type
        
        if best_type is None or best_score == 0:
            best_type = "spreadsheet"
        
        # Calculate confidence
        if best_score >= 3:
            confidence = 0.85
        elif best_score >= 2:
            confidence = 0.75
        elif best_score >= 1:
            confidence = 0.65
        else:
            best_type = "spreadsheet"
            confidence = 0.55  # Slightly better than fallback
        
        return (best_type, confidence, {"matches": best_score, "header_keywords": header_row[:5], "score_details": scores})
    
    except Exception as e:
        return ("spreadsheet", 0.5, {"error": str(e)})


def analyze_cad_document(file_path: str, max_pages: int = 3) -> Dict[str, Any]:
    """Complete CAD analysis: extract metadata, BOM, dimensions."""
    if not file_path.lower().endswith('.pdf'):
        return {'is_cad': False, 'error': 'Only PDF files supported for CAD analysis'}
    
    # Extract text
    text = _extract_pdf_text(file_path, max_pages)
    
    # Check if it's a CAD document
    if not is_cad_document(text):
        return {
            'is_cad': False,
            'file': file_path,
        }
    
    # Extract CAD-specific data
    return {
        'is_cad': True,
        'file': file_path,
        'extracted_text': text[:2000],  # First 2000 chars
        'metadata': extract_cad_metadata(text),
        'bom': extract_cad_bom(text),
        'dimensions': extract_cad_dimensions(text),
    }


