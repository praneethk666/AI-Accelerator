import os
import uuid
import json
from typing import List, Dict, Any, Optional
from pptx import Presentation

def extract_ppt(file_path: str, document_id: Optional[str] = None, config: Optional[Dict[str, Any]] = None) -> List[Dict[str, Any]]:
    """
    Reads a PowerPoint file, extracts slide text and speaker notes,
    and returns data strictly matching the v8.3 NormalizedBlock contract.
    Utilizes a config dictionary to avoid hardcoded metadata.
    """
    blocks: List[Dict[str, Any]] = []
    cfg = config or {}
    
    doc_id = str(document_id) if document_id else str(uuid.uuid4())
    filename = os.path.basename(file_path)
    
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"Failed to read file {file_path}: {e}")
        return blocks
    
    for slide_index, slide in enumerate(prs.slides):
        try:
            slide_num = slide_index + 1
            slide_text_elements = []
            
            # Extract standard text
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_elements.append(shape.text.strip())
                    
            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_text_elements.append(f"Speaker Notes: {notes}")
                    
            full_text = "\n".join(slide_text_elements).strip()
            
            if not full_text:
                continue 

            #8.3 SCHEMA
            block = {
                "block_id": str(uuid.uuid4()),
                "document_id": doc_id,
                "type": "text",              
                "text": full_text,           
                "table_data": None,          
                "source_ref": {
                    "filename": filename,    
                    "page": None,            
                    "sheet": None,           
                    "slide": slide_num,      
                    "bbox": None             
                },
                "confidence": cfg.get("extraction_confidence", 1.0),           
                "language": cfg.get("default_language", "en"),            
                "metadata": {
                    "enrichment_failed": cfg.get("enrichment_failed_flag", False) 
                }
            }
            
            blocks.append(block)
            
        except Exception as e:
            # Fail gracefully
            print(f"Skipping slide {slide_index + 1} due to error: {e}")
            continue
            
    return blocks

# ---------------------------------------------------------
# SANDBOX TEST 
# ---------------------------------------------------------
if __name__ == "__main__":
    test_file_path = "test-data/test.pptx"  
    
    # Simulation of a pipeline config being passed in
    mock_pipeline_config = {
        "extraction_confidence": 0.95,
        "default_language": "en"
    }
    
    print(f"Testing STRICT 8.3 extraction on {test_file_path}...\n")
    if os.path.exists(test_file_path):
        results = extract_ppt(test_file_path, config=mock_pipeline_config)
        
        if results:
            print(f"✅ Success! Extracted {len(results)} slide(s).\n")
            print("--- Output Preview ---")
            print(json.dumps(results[0], indent=2))
        else:
            print("❌ No text extracted.")
    else:
        print(f"⚠️ Dummy file '{test_file_path}' not found. Check test-data folder!")