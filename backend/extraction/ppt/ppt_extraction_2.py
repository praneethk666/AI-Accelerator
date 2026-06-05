import os
import uuid
import json
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Emu


class PPTExtractorTool:
    """
    Extracts text, notes, tables, and embedded images/charts from PowerPoint files.
    Output strictly follows the NormalizedBlock contract in schemas.py.
    """

    def run(self, state: Dict[str, Any], config: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        Entry point for the pipeline.

        state  — expects state["file_path"] and optionally state["document_id"]
        config — pipeline config dict (extraction_confidence, default_language, etc.)
        """
        file_path   = state["file_path"]
        document_id = state.get("document_id")
        return self._extract(file_path, document_id=document_id, config=config)

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    def _extract(
        self,
        file_path: str,
        document_id: Optional[str] = None,
        config: Optional[Dict[str, Any]] = None,
    ) -> List[Dict[str, Any]]:
        blocks: List[Dict[str, Any]] = []
        cfg      = config or {}
        doc_id   = str(document_id) if document_id else str(uuid.uuid4())
        filename = os.path.basename(file_path)

        try:
            prs = Presentation(file_path)
        except Exception as e:
            print(f"[ppt_extractor] Failed to open {file_path}: {e}")
            return blocks

        for slide_index, slide in enumerate(prs.slides):
            slide_num = slide_index + 1
            try:
                # ── 1. Text + notes block per slide ───────────────────
                text_block = self._extract_slide_text(
                    slide, slide_num, doc_id, filename, cfg
                )
                if text_block:
                    blocks.append(text_block)

                # ── 2. Table blocks ────────────────────────────────────
                table_blocks = self._extract_slide_tables(
                    slide, slide_num, doc_id, filename, cfg
                )
                blocks.extend(table_blocks)

                # ── 3. Image / chart blocks ────────────────────────────
                image_blocks = self._extract_slide_images(
                    slide, slide_num, doc_id, filename, cfg
                )
                blocks.extend(image_blocks)

            except Exception as e:
                print(f"[ppt_extractor] Skipping slide {slide_num}: {e}")
                continue

        return blocks

    def _extract_slide_text(
        self,
        slide,
        slide_num: int,
        doc_id: str,
        filename: str,
        cfg: Dict[str, Any],
    ) -> Optional[Dict[str, Any]]:
        """One text block per slide — all text shapes + speaker notes combined."""
        parts = []

        for shape in slide.shapes:
            # Skip tables and charts — handled separately
            if shape.has_table or shape.shape_type == 3:  # 3 = MSO_SHAPE_TYPE.CHART
                continue
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"Speaker Notes: {notes}")

        full_text = "\n".join(parts).strip()
        if not full_text:
            return None

        return {
            "block_id":    str(uuid.uuid4()),
            "document_id": doc_id,
            "type":        "text",
            "text":        full_text,
            "table_data":  None,
            "source_ref": {
                "filename": filename,
                "page":     None,
                "sheet":    None,
                "slide":    slide_num,
                "bbox":     None,
            },
            "confidence": cfg.get("extraction_confidence", 1.0),
            "language":   cfg.get("default_language", "en"),
            "metadata": {
                "enrichment_failed": cfg.get("enrichment_failed_flag", False),
            },
        }

    def _extract_slide_tables(
        self,
        slide,
        slide_num: int,
        doc_id: str,
        filename: str,
        cfg: Dict[str, Any],
    ) -> List[Dict[str, Any]]:
        """One table block per table shape on the slide."""
        import pandas as pd

        blocks = []

        for shape in slide.shapes:
            if not shape.has_table:
                continue
            try:
                tbl     = shape.table
                rows    = []
                headers = []

                for i, row in enumerate(tbl.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    if i == 0:
                        headers = cells
                    else:
                        rows.append(cells)

                df       = pd.DataFrame(rows, columns=headers if headers else None)
                block_id = str(uuid.uuid4())

                blocks.append({
                    "block_id":    block_id,
                    "document_id": doc_id,
                    "type":        "table",
                    "text":        df.to_markdown(index=False),
                    "table_data": {
                        "headers": headers,
                        "rows":    rows,
                    },
                    "source_ref": {
                        "filename": filename,
                        "page":     None,
                        "sheet":    None,
                        "slide":    slide_num,
                        "bbox":     None,
                    },
                    "confidence": cfg.get("extraction_confidence", 1.0),
                    "language":   cfg.get("default_language", "en"),
                    "metadata": {
                        "enrichment_failed": cfg.get("enrichment_failed_flag", False),
                    },
                })

            except Exception as e:
                print(f"[ppt_extractor] Skipping table on slide {slide_num}: {e}")
                continue

        return blocks

    def _extract_slide_images(
        self,
        slide,
        slide_num: int,
        doc_id: str,
        filename: str,
        cfg: Dict[str, Any],
    ) -> List[Dict[str, Any]]:
        """
        Saves embedded image/chart bytes to uploads/images/<doc_id>/<block_id>_raw.png.
        Returns image blocks with metadata["raw_image_path"] set.
        vision_enrichment_tool picks them up from there and fills text + image_path.
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        blocks  = []
        out_dir = os.path.join("uploads", "images", doc_id)
        os.makedirs(out_dir, exist_ok=True)

        for shape in slide.shapes:
            image_bytes = None
            shape_label = None

            try:
                # ── Inline pictures ────────────────────────────────────
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_bytes = shape.image.blob
                    shape_label = getattr(shape, "name", "picture")

                # ── Embedded charts ────────────────────────────────────
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    # Charts don't expose raw bytes via python-pptx directly.
                    # We grab the chart's embedded image if it exists (cached PNG
                    # that PowerPoint stores for non-Excel renderers).
                    chart_part = shape.chart._part
                    for rel in chart_part.rels.values():
                        if "image" in rel.reltype:
                            image_bytes = rel.target_part.blob
                            break
                    shape_label = getattr(shape, "name", "chart")

                if not image_bytes:
                    continue

                block_id = str(uuid.uuid4())
                raw_path = os.path.join(out_dir, f"{block_id}_raw.png")

                with open(raw_path, "wb") as f:
                    f.write(image_bytes)

                blocks.append({
                    "block_id":    block_id,
                    "document_id": doc_id,
                    "type":        "image",
                    # text is empty — vision_enrichment_tool fills this in
                    "text":        None,
                    "table_data":  None,
                    "source_ref": {
                        "filename": filename,
                        "page":     None,
                        "sheet":    None,
                        "slide":    slide_num,
                        "bbox":     None,
                    },
                    "confidence": cfg.get("extraction_confidence", 1.0),
                    "language":   cfg.get("default_language", "en"),
                    "metadata": {
                        # extractor sets raw_image_path; vision sets image_path + text
                        "raw_image_path":   raw_path,
                        "shape_label":      shape_label,
                        "enrichment_failed": False,
                    },
                })

            except Exception as e:
                print(f"[ppt_extractor] Skipping image/chart on slide {slide_num}: {e}")
                continue

        return blocks


# ------------------------------------------------------------------
# SANDBOX TEST
# ------------------------------------------------------------------
if __name__ == "__main__":
    test_file = "test-data/test.pptx"

    mock_state  = {"file_path": test_file, "document_id": "doc-002"}
    mock_config = {"extraction_confidence": 0.95, "default_language": "en"}

    tool    = PPTExtractorTool()
    results = tool.run(mock_state, mock_config)

    texts   = [b for b in results if b["type"] == "text"]
    tables  = [b for b in results if b["type"] == "table"]
    images  = [b for b in results if b["type"] == "image"]

    print(f"Extracted {len(texts)} text block(s), {len(tables)} table(s), {len(images)} image/chart block(s).\n")

    if texts:
        print("--- First text block ---")
        print(json.dumps(texts[0], indent=2))

    if tables:
        print("\n--- First table block ---")
        preview = {k: v for k, v in tables[0].items() if k != "table_data"}
        preview["text_preview"] = (tables[0]["text"] or "")[:300]
        print(json.dumps(preview, indent=2))

    if images:
        print("\n--- First image block ---")
        print(json.dumps(images[0], indent=2))