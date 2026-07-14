import os
import uuid
import re
import zipfile
from io import BytesIO
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from backend.core.paths import display_filename
from backend.core.schemas import NormalizedBlock, SourceRef, as_dicts


def _detect_image_ext(data: bytes) -> str:
    """Sniff an embedded image's real format from its magic bytes."""
    if not data:
        return "png"
    if data[:8] == b"\x89PNG\r\n\x1a\n": return "png"
    if data[:3] == b"\xff\xd8\xff": return "jpg"
    if data[:6] in (b"GIF87a", b"GIF89a"): return "gif"
    if data[:4] == b"RIFF" and data[8:12] == b"WEBP": return "webp"
    if data[:4] in (b"MM\x00*", b"II*\x00"): return "tiff"
    if data[:2] == b"BM": return "bmp"
    if data[:2] == b"\x1f\x8b": return "gz"
    return "png"


def _safe_text(value) -> str:
    if value is None:
        return ""
    try:
        return str(value).strip()
    except Exception:
        return ""


def _save_image_blob(image_bytes: bytes, out_dir: str, block_id: str) -> tuple[str, str]:
    """Save raw embedded bytes and, if possible, a PNG preview."""
    os.makedirs(out_dir, exist_ok=True)
    ext = _detect_image_ext(image_bytes)
    raw_path = os.path.join(out_dir, f"{block_id}_raw.{ext}")
    with open(raw_path, "wb") as f:
        f.write(image_bytes)

    preview_path = raw_path
    try:
        from PIL import Image

        img = Image.open(BytesIO(image_bytes))
        img = img.convert("RGBA") if img.mode not in ("RGB", "RGBA") else img
        preview_path = os.path.join(out_dir, f"{block_id}.png")
        img.save(preview_path, format="PNG")
    except Exception:
        preview_path = raw_path

    return raw_path, preview_path


class PPTExtractorTool:
    """
    Extracts text, notes, tables, and embedded images/charts from PowerPoint files.
    Output strictly follows the NormalizedBlock contract in schemas.py.

    """
    name = "ppt_extraction"

    @staticmethod
    def _normalize_file_type(state: dict) -> None:
        """Align categorize output with pipeline dispatch.

        The classifier may label PPTs as 'powerpoint', but the graph dispatch
        expects 'ppt'. Normalize early so extraction actually runs.
        """
        if state.get("file_type") == "powerpoint":
            state["file_type"] = "ppt"

    @staticmethod
    def _unsupported_legacy_ppt(state: dict, file_path: str) -> bool:
        """python-pptx cannot read binary `.ppt`; fail softly and visibly."""
        if os.path.splitext(file_path)[1].lower() != ".ppt":
            return False
        state.setdefault("errors", []).append(
            "ppt_extraction: legacy .ppt files are not supported by python-pptx; convert to .pptx"
        )
        state.setdefault("blocks", [])
        return True

    def _fallback_from_pptx_xml(self, file_path: str, document_id: str, config: dict) -> list[NormalizedBlock]:
        """Try to salvage text directly from PPTX XML when python-pptx fails."""
        if os.path.splitext(file_path)[1].lower() != ".pptx":
            return []
        try:
            with zipfile.ZipFile(file_path) as zf:
                slide_names = sorted(
                    name for name in zf.namelist()
                    if re.match(r"ppt/slides/slide\d+\.xml$", name)
                )
                blocks: list[NormalizedBlock] = []
                for slide_index, slide_name in enumerate(slide_names, start=1):
                    try:
                        raw = zf.read(slide_name).decode("utf-8", errors="ignore")
                    except Exception:
                        continue
                    text = re.sub(r"<[^>]+>", " ", raw)
                    text = re.sub(r"\s+", " ", text).strip()
                    if not text:
                        continue
                    blocks.append(NormalizedBlock(
                        block_id=str(uuid.uuid4()),
                        document_id=document_id,
                        type="text",
                        text=text[:5000],
                        source_ref=SourceRef(
                            filename=display_filename(file_path),
                            slide=slide_index,
                        ),
                        confidence=0.2,
                        language="en",
                        metadata={
                            "fallback": True,
                            "source": "pptx_xml",
                            "enrichment_failed": config.get("enrichment_failed_flag", False),
                        },
                    ))
                return blocks
        except Exception:
            return []

    def run(self, state: dict, config: dict) -> dict:
        self._normalize_file_type(state)
        file_path   = state["file_path"]
        document_id = state.get("document_id")
        if self._unsupported_legacy_ppt(state, file_path):
            return state
        blocks      = self._extract(file_path, document_id=document_id, config=config)
        if not blocks:
            blocks = self._fallback_from_pptx_xml(file_path, str(document_id or uuid.uuid4()), config)
        state["blocks"] = as_dicts(blocks)
        state.setdefault("errors", [])
        return state

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    def _iter_shapes(self, shapes):
        """Recursively yield shapes, descending into grouped shapes (plain
        slide.shapes iteration misses everything nested inside a group)."""
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from self._iter_shapes(shape.shapes)
            else:
                yield shape

    def _detect_language(self, text: str, default: str = "en") -> str:
        """Best-effort ISO 639-1 language code; falls back to `default`."""
        if not text or not text.strip():
            return default
        try:
            from langdetect import detect
            return detect(text)
        except Exception:
            return default

    def _extract(
        self,
        file_path: str,
        document_id: Optional[str] = None,
        config: Optional[Dict[str, Any]] = None,
    ) -> List[NormalizedBlock]:
        blocks: List[NormalizedBlock] = []
        cfg      = config or {}
        doc_id   = str(document_id) if document_id else str(uuid.uuid4())
        filename = display_filename(file_path)

        try:
            prs = Presentation(file_path)
        except Exception:
            return blocks

        for slide_index, slide in enumerate(prs.slides):
            slide_num = slide_index + 1
            try:
                text_block = self._extract_slide_text(slide, slide_num, doc_id, filename, cfg)
                if text_block:
                    blocks.append(text_block)

                blocks.extend(self._extract_slide_tables(slide, slide_num, doc_id, filename, cfg))
                blocks.extend(self._extract_slide_images(slide, slide_num, doc_id, filename, cfg))

            except Exception:
                continue

        return blocks

    def _extract_slide_text(
        self,
        slide,
        slide_num: int,
        doc_id: str,
        filename: str,
        cfg: Dict[str, Any],
    ) -> Optional[NormalizedBlock]:
        parts = []

        for shape in self._iter_shapes(slide.shapes):
            if shape.has_table or shape.shape_type == 3:
                continue
            if hasattr(shape, "text"):
                text = _safe_text(getattr(shape, "text", ""))
                if text:
                    parts.append(text)

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = _safe_text(slide.notes_slide.notes_text_frame.text)
            if notes:
                parts.append(f"Speaker Notes: {notes}")

        full_text = "\n".join(parts).strip()
        if not full_text:
            return None

        return NormalizedBlock(
            block_id=str(uuid.uuid4()),
            document_id=doc_id,
            type="text",
            text=full_text,
            source_ref=SourceRef(
                filename=filename,
                slide=slide_num,
            ),
            confidence=cfg.get("extraction_confidence", 1.0),
            language=self._detect_language(full_text, cfg.get("default_language", "en")),
            metadata={
                "enrichment_failed": cfg.get("enrichment_failed_flag", False),
            },
        )

    def _extract_slide_tables(
        self,
        slide,
        slide_num: int,
        doc_id: str,
        filename: str,
        cfg: Dict[str, Any],
    ) -> List[NormalizedBlock]:
        import pandas as pd
        blocks = []

        for shape in self._iter_shapes(slide.shapes):
            if not shape.has_table:
                continue
            try:
                tbl     = shape.table
                rows    = []
                headers = []

                for i, row in enumerate(tbl.rows):
                    cells = [_safe_text(cell.text) for cell in row.cells]
                    if i == 0:
                        headers = cells
                    else:
                        rows.append(cells)

                df       = pd.DataFrame(rows, columns=headers if headers else None)
                block_id = str(uuid.uuid4())
                md_text  = df.to_markdown(index=False)

                blocks.append(NormalizedBlock(
                    block_id=block_id,
                    document_id=doc_id,
                    type="table",
                    text=md_text,
                    table_data={
                        "headers": headers,
                        "rows":    rows,
                    },
                    source_ref=SourceRef(
                        filename=filename,
                        slide=slide_num,
                    ),
                    confidence=cfg.get("extraction_confidence", 1.0),
                    language=self._detect_language(md_text, cfg.get("default_language", "en")),
                    metadata={
                        "enrichment_failed": cfg.get("enrichment_failed_flag", False),
                    },
                ))

            except Exception:
                continue

        return blocks

    def _extract_slide_images(
        self,
        slide,
        slide_num: int,
        doc_id: str,
        filename: str,
        cfg: Dict[str, Any],
    ) -> List[NormalizedBlock]:
        blocks  = []
        out_dir = os.path.join("uploads", "images", doc_id)
        os.makedirs(out_dir, exist_ok=True)

        for shape in self._iter_shapes(slide.shapes):
            image_bytes = None
            shape_label = None

            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_bytes = shape.image.blob
                    shape_label = _safe_text(getattr(shape, "name", "picture")) or "picture"

                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    chart_part = shape.chart._part
                    for rel in chart_part.rels.values():
                        if "image" in rel.reltype:
                            image_bytes = rel.target_part.blob
                            break
                    shape_label = _safe_text(getattr(shape, "name", "chart")) or "chart"

                if not image_bytes:
                    continue

                block_id = str(uuid.uuid4())
                raw_path, preview_path = _save_image_blob(image_bytes, out_dir, block_id)

                blocks.append(NormalizedBlock(
                    block_id=block_id,
                    document_id=doc_id,
                    type="image_caption",
                    text="",
                    source_ref=SourceRef(
                        filename=filename,
                        slide=slide_num,
                    ),
                    confidence=cfg.get("extraction_confidence", 1.0),
                    language=cfg.get("default_language", "en"),
                    metadata={
                        "raw_image_path":    raw_path,
                        "image_path":        preview_path if preview_path else raw_path,
                        "pending_vision":    True,
                        "shape_label":       shape_label,
                        "enrichment_failed": False,
                    },
                ))

            except Exception:
                continue

        return blocks


# ------------------------------------------------------------------
# SANDBOX TEST
# ------------------------------------------------------------------
if __name__ == "__main__":
    test_file = "test-data/test.pptx"

    mock_state  = {"file_path": test_file, "document_id": "doc-002"}
    mock_config = {"extraction_confidence": 0.95, "default_language": "en"}

    tool    = PPTExtractorTool()
    state   = tool.run(mock_state, mock_config)
    results = state["blocks"]

    texts   = [b for b in results if b.type == "text"]
    tables  = [b for b in results if b.type == "table"]
    images  = [b for b in results if b.type == "image_caption"]

    print(f"Extracted {len(texts)} text block(s), {len(tables)} table(s), {len(images)} image/chart block(s).\n")

    if texts:
        print("--- First text block ---")
        print(f"  text preview: {texts[0].text[:200]}")

    if tables:
        print("\n--- First table block ---")
        print(f"  text preview: {(tables[0].text or '')[:200]}")
        print(f"  headers: {tables[0].table_data['headers']}")

    if images:
        print("\n--- First image block ---")
        print(f"  raw_image_path: {images[0].metadata['raw_image_path']}")
        print(f"  pending_vision: {images[0].metadata['pending_vision']}")
